import os
import sys
from pptx import Presentation
from PIL import Image
import fitz  # PyMuPDF

# Function to create PowerPoint from PDF
def pdf_to_pptx(pdf_path, title, recipient):
    presentation = Presentation()
    
    # Open the PDF
    doc = fitz.open(pdf_path)

    # Loop through each page of the PDF
    for page_num in range(doc.page_count):
        page = doc.load_page(page_num)  # Load the page
        text = page.get_text("text")   # Extract text (optional for titles/captions)
        
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide layout
        
        # Set the title
        title_box = slide.shapes.title
        title_box.text = f"{title} - Page {page_num + 1}"

        # Add the page as an image (optional step to render PDF as an image)
        pix = page.get_pixmap()  # Render page as an image
        img_path = f"uploads/page_{page_num}.png"
        pix.save(img_path)

        # Add image to slide
        slide.shapes.add_picture(img_path, 0, 0, width=presentation.slide_width, height=presentation.slide_height)
        
        # Clean up the image
        os.remove(img_path)

    # Save the presentation
    output_file = f"docs/output/{title.replace(' ', '_')}_presentation.pptx"
    presentation.save(output_file)
    print(f"Presentation saved at {output_file}")
    
    return
