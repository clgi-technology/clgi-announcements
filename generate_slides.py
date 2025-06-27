import os
import sys
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from pptx import Presentation
from PIL import Image
import fitz  # PyMuPDF

# Function to create PowerPoint from PDF
def pdf_to_pptx(pdf_path, title, recipient):
    presentation = Presentation()
    
    # Open the PDF
    doc = fitz.open(pdf_path)

    # Loop through each page of the PDF
    for page_num in range(doc.page_count):
        page = doc.load_page(page_num)  # Load the page
        text = page.get_text("text")   # Extract text (optional for titles/captions)
        
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide layout
        
        # Set the title
        title_box = slide.shapes.title
        title_box.text = f"{title} - Page {page_num + 1}"

        # Add the page as an image (optional step to render PDF as an image)
        pix = page.get_pixmap()  # Render page as an image
        img_path = f"uploads/page_{page_num}.png"
        pix.save(img_path)

        # Add image to slide
        slide.shapes.add_picture(img_path, 0, 0, width=presentation.slide_width, height=presentation.slide_height)
        
        # Clean up the image
        os.remove(img_path)

    # Save the presentation
    output_file = f"docs/output/{title.replace(' ', '_')}_presentation.pptx"
    presentation.save(output_file)
    print(f"Presentation saved at {output_file}")
    
    return output_file

# Function to handle JPEG to PPTX
def jpeg_to_pptx(jpeg_path, title, recipient):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide layout
    
    # Add the JPEG as a picture to the slide
    slide.shapes.add_picture(jpeg_path, 0, 0, width=presentation.slide_width, height=presentation.slide_height)
    
    # Save the presentation
    output_file = f"docs/output/{title.replace(' ', '_')}_presentation.pptx"
    presentation.save(output_file)
    print(f"JPEG converted and saved at {output_file}")
    
    return output_file

# Function to send email notification
def send_email(subject, body, recipient_email):
    sender_email = os.getenv("SENDER_EMAIL")  # Set this in your environment variables
    sender_password = os.getenv("SENDER_PASSWORD")  # Set this in your environment variables

    # Setup the MIME
    msg = MIMEMultipart()
    msg['From'] = sender_email
    msg['To'] = recipient_email
    msg['Subject'] = subject

    msg.attach(MIMEText(body, 'plain'))

    try:
        # Connect to Gmail's SMTP server and send the email
        server = smtplib.SMTP('smtp.gmail.com', 587)
        server.starttls()  # Secure connection
        server.login(sender_email, sender_password)
        text = msg.as_string()
        server.sendmail(sender_email, recipient_email, text)
        print(f"Email sent to {recipient_email}")
    except Exception as e:
        print(f"Failed to send email: {str(e)}")
    finally:
        server.quit()

# Main function to convert files
def main(file_path, title, recipient):
    file_extension = file_path.split('.')[-1].lower()
    
    if file_extension == 'pdf':
        output_file = pdf_to_pptx(file_path, title, recipient)
    elif file_extension in ['jpg', 'jpeg']:
        output_file = jpeg_to_pptx(file_path, title, recipient)
    else:
        print("Unsupported file type")
        sys.exit(1)

    # Send email notification
    subject = f"Your PowerPoint Presentation: {title}"
    body = f"Your PowerPoint presentation has been created successfully.\n\nDownload link: {output_file}"
    send_email(subject, body, recipient)
    
    return output_file

if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: python generate_slides.py <file_path> <title> <recipient>")
        sys.exit(1)
    
    file_path = sys.argv[1]
    title = sys.argv[2]
    recipient = sys.argv[3]
    
    main(file_path, title, recipient)
